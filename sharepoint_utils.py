import os
import json
import logging
import requests
import msal
from typing import Dict, List, Any, Optional
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

logging.basicConfig(level=logging.INFO,
                   format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


def normalize_path(path):
    """標準化路徑格式，便於比較"""
    if not path:
        return ""
    # 移除開頭和結尾的斜線，統一使用正斜線
    return path.strip("/\\").replace("\\", "/")

# auth  and  basic
class SharePointClient:
    """SharePoint 客戶端，處理認證和基本操作"""

    def __init__(self, client_id=None, client_secret=None, tenant_id=None,
                 site_name=None, tenant_name=None, config_path=None):
        """初始化 SharePoint 客戶端

        可以通過直接提供參數或配置文件初始化

        Args:
            client_id (str, optional): Azure AD 應用程式的 Client ID
            client_secret (str, optional): Azure AD 應用程式的 Client Secret
            tenant_id (str, optional): Azure AD 租戶 ID
            site_name (str, optional): SharePoint 網站名稱
            tenant_name (str, optional): SharePoint 租戶名稱
            config_path (str, optional): 配置文件路徑
        """
        if config_path:
            self.config = self._load_config(config_path)
            self.client_id = self.config.get("client_id")
            self.client_secret = self.config.get("secret")
            self.tenant_id = self.config.get("tenant_id")
            self.site_name = self.config.get("site_name")
            self.tenant_name = self.config.get("tenant")
            self.scope = self.config.get("scope", ["https://graph.microsoft.com/.default"])
        else:
            self.client_id = client_id
            self.client_secret = client_secret
            self.tenant_id = tenant_id
            self.site_name = site_name
            self.tenant_name = tenant_name
            self.scope = ["https://graph.microsoft.com/.default"]

        self.base_url = "https://graph.microsoft.com/v1.0"
        self.access_token = None
        self.drive_id = None

    def _load_config(self, config_path):
        """載入配置文件"""
        try:
            with open(config_path) as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"載入配置文件失敗: {str(e)}")
            raise

    def get_access_token(self):
        """獲取訪問令牌"""
        if self.access_token:
            return self.access_token

        try:
            # 確保使用正確的 authority URL 格式
            authority = f"https://login.microsoftonline.com/{self.tenant_id}"

            app = msal.ConfidentialClientApplication(
                self.client_id,
                authority=authority,
                client_credential=self.client_secret,
            )

            result = app.acquire_token_silent(self.scope, account=None)
            if not result:
                result = app.acquire_token_for_client(scopes=self.scope)

            if "access_token" in result:
                self.access_token = result["access_token"]
                logger.info("✅ 成功獲取訪問令牌")
                return self.access_token
            else:
                error_msg = f"❌ 無法獲取訪問令牌: {result.get('error_description', '未知錯誤')}"
                logger.error(error_msg)
                return None

        except Exception as e:
            logger.error(f"❌ 獲取訪問令牌時發生錯誤: {str(e)}")
            return None

    def get_site_id(self):
        """獲取 SharePoint 站點 ID"""
        try:
            # 確保有訪問令牌
            if not self.access_token:
                self.access_token = self.get_access_token()
                if not self.access_token:
                    return None

            # 獲取站點ID
            site_url = f"{self.base_url}/sites/{self.tenant_name}.sharepoint.com:/sites/{self.site_name}"
            response = requests.get(
                site_url,
                headers={'Authorization': f'Bearer {self.access_token}'}
            )

            if response.status_code != 200:
                logger.error(f"❌ 無法獲取站點信息: {response.text}")
                return None

            site_id = response.json()['id']
            logger.info(f"✅ 成功獲取站點 ID: {site_id}")
            return site_id

        except Exception as e:
            logger.error(f"❌ 獲取站點 ID 時發生錯誤: {str(e)}")
            return None

    def get_drive_id(self):
        """獲取 SharePoint 驅動器 ID"""
        if self.drive_id:
            return self.drive_id

        try:
            # 確保有訪問令牌
            if not self.access_token:
                self.access_token = self.get_access_token()
                if not self.access_token:
                    return None

            # 1. 首先獲取站點ID
            site_id = self.get_site_id()
            if not site_id:
                return None

            # 2. 使用站點ID獲取驅動器信息
            drives_url = f"{self.base_url}/sites/{site_id}/drives"
            response = requests.get(
                drives_url,
                headers={'Authorization': f'Bearer {self.access_token}'}
            )

            if response.status_code != 200:
                logger.error(f"❌ 無法獲取驅動器信息: {response.text}")
                return None

            # 通常第一個驅動器是文檔庫
            self.drive_id = response.json()['value'][0]['id']
            logger.info(f"✅ 成功獲取驅動器 ID: {self.drive_id}")
            return self.drive_id

        except Exception as e:
            logger.error(f"❌ 獲取驅動器 ID 時發生錯誤: {str(e)}")
            return None

    def get_items_recursive(self, item_id=None, current_path=""):
        """遞歸獲取資料夾和文件的資訊

        Args:
            item_id (str, optional): 項目 ID，如果為 None 則獲取根目錄
            current_path (str): 當前路徑

        Returns:
            list: 項目列表
        """
        items_info = []

        try:
            # 確保有訪問令牌和驅動器ID
            if not self.access_token:
                self.access_token = self.get_access_token()
                if not self.access_token:
                    return items_info

            if not self.drive_id:
                self.drive_id = self.get_drive_id()
                if not self.drive_id:
                    return items_info

            # 構建 URL
            if not item_id:
                url = f"{self.base_url}/drives/{self.drive_id}/root/children"
            else:
                url = f"{self.base_url}/drives/{self.drive_id}/items/{item_id}/children"

            # 發送請求
            response = requests.get(
                url,
                headers={'Authorization': f'Bearer {self.access_token}'}
            )

            if response.status_code != 200:
                logger.error(f"❌ 獲取項目列表失敗: {response.text}")
                return items_info

            items = response.json().get('value', [])

            # 處理每個項目
            for item in items:
                name = item.get('name', '')
                item_id = item.get('id', '')
                is_folder = 'folder' in item

                # 構建項目路徑
                item_path = f"{current_path}/{name}" if current_path else name

                # 創建項目信息
                item_info = {
                    'name': name,
                    'path': item_path,
                    'type': 'folder' if is_folder else 'file',
                    'id': item_id,
                    'created_time': item.get('createdDateTime'),
                    'modified_time': item.get('lastModifiedDateTime'),
                    'size': item.get('size', 0)
                }

                # 如果是文件，添加下載 URL
                if not is_folder and 'file' in item:
                    item_info['download_url'] = item.get('@microsoft.graph.downloadUrl')

                items_info.append(item_info)

                # 如果是資料夾，遞歸獲取子項目
                if is_folder:
                    children = self.get_items_recursive(item_id, item_path)
                    items_info.extend(children)

            return items_info

        except Exception as e:
            logger.error(f"❌ 獲取項目列表時發生錯誤: {str(e)}")
            return items_info

    def get_list_item_fields(self, item_id=None, path=""):
        """遞歸獲取SharePoint文件和資料夾的所有欄位值"""
        fields_list = []

        try:
            # 確保有訪問令牌和驅動器ID
            if not self.access_token:
                self.access_token = self.get_access_token()
                if not self.access_token:
                    return fields_list

            if not self.drive_id:
                self.drive_id = self.get_drive_id()
                if not self.drive_id:
                    return fields_list

            # 獲取根目錄或指定資料夾的項目
            if not item_id:
                url = f"{self.base_url}/drives/{self.drive_id}/root/children"
            else:
                url = f"{self.base_url}/drives/{self.drive_id}/items/{item_id}/children"

            res = requests.get(url, headers={'Authorization': f'Bearer {self.access_token}'})
            res.raise_for_status()  # 如果請求失敗會拋出異常
            items = res.json().get("value", [])

            for item in items:
                item_id = item.get("id")
                name = item.get("name")
                current_path = f"{path}/{name}" if path else name
                is_folder = 'folder' in item

                # 構建基本項目信息
                item_info = {
                    "name": current_path.replace("\\", "/"),
                    "type": "folder" if is_folder else "file",
                    "path": current_path,
                    "created_time": item.get("createdDateTime"),
                    "modified_time": item.get("lastModifiedDateTime"),
                    "size": item.get("size", 0)
                }

                # 如果是文件，獲取額外的欄位信息
                if not is_folder:
                    fields_url = f"{self.base_url}/drives/{self.drive_id}/items/{item_id}/listItem/fields"
                    fields_res = requests.get(fields_url, headers={'Authorization': f'Bearer {self.access_token}'})

                    if fields_res.status_code == 200:
                        fields = fields_res.json()
                        # 添加所有自定義欄位
                        for k, v in fields.items():
                            if k not in ['id', 'FileLeafRef'] and not k.startswith('_'):
                                item_info[k] = v

                fields_list.append(item_info)

                # 如果是資料夾，遞歸處理
                if is_folder:
                    child_fields = self.get_list_item_fields(
                        item_id=item_id,
                        path=current_path
                    )
                    fields_list.extend(child_fields)

            return fields_list

        except requests.exceptions.RequestException as e:
            logger.error(f"API請求錯誤: {str(e)}")
            return []
        except Exception as e:
            logger.error(f"處理文件時發生錯誤: {str(e)}")
            return []

    def download_file(self, file_path, local_path=None):
        """下載文件

        Args:
            file_path (str): SharePoint中的文件路徑
            local_path (str, optional): 本地保存路徑，如果為None則返回文件內容

        Returns:
            bool或bytes: 如果指定了local_path，則返回下載是否成功；否則返回文件內容
        """
        try:
            # 確保有訪問令牌和驅動器ID
            if not self.access_token:
                self.access_token = self.get_access_token()
                if not self.access_token:
                    return False if local_path else None

            if not self.drive_id:
                self.drive_id = self.get_drive_id()
                if not self.drive_id:
                    return False if local_path else None

            # 構建下載URL
            download_url = f"{self.base_url}/drives/{self.drive_id}/root:/{file_path}:/content"

            # 發送請求下載文件
            response = requests.get(
                download_url,
                headers={'Authorization': f'Bearer {self.access_token}'},
                stream=True
            )

            if response.status_code == 200:
                if local_path:
                    # 確保目錄存在
                    os.makedirs(os.path.dirname(local_path), exist_ok=True)

                    # 保存文件
                    with open(local_path, 'wb') as f:
                        for chunk in response.iter_content(chunk_size=8192):
                            f.write(chunk)
                    return True
                else:
                    return response.content
            else:
                logger.error(f"❌ 下載文件失敗: {response.status_code} - {response.text}")
                return False if local_path else None
        except Exception as e:
            logger.error(f"❌ 下載文件時發生錯誤: {str(e)}")
            return False if local_path else None

    def create_folder(self, folder_path):
        """在 SharePoint 中建立資料夾"""
        try:
            # 確保有訪問令牌和驅動器ID
            if not self.access_token:
                self.access_token = self.get_access_token()
                if not self.access_token:
                    return None

            if not self.drive_id:
                self.drive_id = self.get_drive_id()
                if not self.drive_id:
                    return None

            url = f"{self.base_url}/drives/{self.drive_id}/root:/{folder_path}"

            # 檢查資料夾是否已存在
            response = requests.get(
                url,
                headers={'Authorization': f'Bearer {self.access_token}'}
            )

            if response.status_code == 200:
                return response.json()['id']

            # 建立新資料夾
            folder_data = {
                "name": os.path.basename(folder_path),
                "folder": {},
                "@microsoft.graph.conflictBehavior": "replace"
            }

            parent_path = os.path.dirname(folder_path)
            if parent_path:
                url = f"{self.base_url}/drives/{self.drive_id}/root:/{parent_path}:/children"
            else:
                url = f"{self.base_url}/drives/{self.drive_id}/root/children"

            response = requests.post(
                url,
                headers={
                    'Authorization': f'Bearer {self.access_token}',
                    'Content-Type': 'application/json'
                },
                json=folder_data
            )

            if response.status_code in [201, 200]:
                return response.json()['id']
            else:
                logger.error(f"❌ 建立資料夾失敗: {response.text}")
                return None
        except Exception as e:
            logger.error(f"❌ 建立資料夾時發生錯誤: {str(e)}")
            return None

    def upload_file(self, file_path, content, metadata=None):
        """上傳檔案到 SharePoint"""
        try:
            # 確保有訪問令牌和驅動器ID
            if not self.access_token:
                self.access_token = self.get_access_token()
                if not self.access_token:
                    return {'success': False, 'error': '無法獲取訪問令牌'}

            if not self.drive_id:
                self.drive_id = self.get_drive_id()
                if not self.drive_id:
                    return {'success': False, 'error': '無法獲取驅動器ID'}

            # 如果檔案小於 4MB，使用單一請求上傳
            if len(content) < 4 * 1024 * 1024:
                url = f"{self.base_url}/drives/{self.drive_id}/root:/{file_path}:/content"
                response = requests.put(
                    url,
                    headers={
                        'Authorization': f'Bearer {self.access_token}',
                        'Content-Type': 'application/octet-stream'
                    },
                    data=content
                )

                if response.status_code in [200, 201]:
                    result = {'success': True, 'id': response.json()['id']}
                else:
                    result = {'success': False, 'error': response.text}

                # 如果有metadata且上傳成功，更新metadata
                if metadata and result['success']:
                    self.update_file_metadata(result['id'], metadata)

                return result
            else:
                # TODO: 實作大檔案的分段上傳
                logger.error("❌ 檔案大小超過4MB，分段上傳尚未實作")
                return {'success': False, 'error': "檔案大小超過4MB，分段上傳尚未實作"}

        except Exception as e:
            logger.error(f"❌ 上傳檔案時發生錯誤: {str(e)}")
            return {'success': False, 'error': str(e)}

    def update_file_metadata(self, item_id, metadata):
        """更新檔案的metadata"""
        try:
            # 確保有訪問令牌和驅動器ID
            if not self.access_token:
                self.access_token = self.get_access_token()
                if not self.access_token:
                    return False

            if not self.drive_id:
                self.drive_id = self.get_drive_id()
                if not self.drive_id:
                    return False

            url = f"{self.base_url}/drives/{self.drive_id}/items/{item_id}/fields"

            response = requests.patch(
                url,
                headers={
                    'Authorization': f'Bearer {self.access_token}',
                    'Content-Type': 'application/json'
                },
                json=metadata
            )

            if response.status_code == 200:
                return True
            else:
                logger.error(f"❌ 更新metadata失敗: {response.text}")
                return False
        except Exception as e:
            logger.error(f"❌ 更新metadata時發生錯誤: {str(e)}")
            return False


# File content extraction functions
def extract_text_from_docx(filepath):
    """Extract text content and keywords from Word document"""
    try:
        # First check if file exists and is readable
        if not os.path.exists(filepath):
            logger.warning(f"Word document not found: {filepath}")
            # Fallback to filename extraction
            filename = os.path.basename(filepath)
            name_without_ext = os.path.splitext(filename)[0]
            result = extract_keywords(name_without_ext)
            return result if result else "📄 " + os.path.splitext(filename)[0]
            
        # Normalize path separators
        normalized_path = filepath.replace('\\', '/')
        
        try:
            doc = Document(normalized_path)
            
            # Extract title (if exists)
            title = ""
            for para in doc.paragraphs:
                if para.style.name.startswith('Heading') and para.text.strip():
                    title = para.text.strip()
                    break

            # Extract text from 2 non-empty paragraphs
            paragraphs_text = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs_text.append(text)
                    if len(paragraphs_text) >= 2:
                        break

            # Combine text
            full_text = " ".join(paragraphs_text)

            # If title exists, add it to the result
            if title:
                result = f"{title} - {extract_keywords(full_text)}"
            else:
                result = extract_keywords(full_text)

            # If no content extracted, fallback to filename
            if not result or result == "No keywords":
                logger.warning(f"No content extracted from document, using filename: {filepath}")
                filename = os.path.basename(filepath)
                name_without_ext = os.path.splitext(filename)[0]
                result = extract_keywords(name_without_ext)
                
            return result if result else "📄 " + os.path.splitext(os.path.basename(filepath))[0]
            
        except Exception as doc_error:
            # Log specific document parsing error
            logger.warning(f"Failed to parse Word document content: {filepath} - {str(doc_error)}")
            # Fallback to filename extraction
            filename = os.path.basename(filepath)
            name_without_ext = os.path.splitext(filename)[0]
            result = extract_keywords(name_without_ext)
            logger.info(f"Using filename keywords for {filepath}: {result}")
            return result if result else "📄 " + name_without_ext
            
    except Exception as e:
        # Final fallback for any other errors
        logger.warning(f"Error processing Word document {filepath}, using filename: {str(e)}")
        try:
            filename = os.path.basename(filepath)
            name_without_ext = os.path.splitext(filename)[0]
            result = extract_keywords(name_without_ext)
            return result if result else "📄 " + name_without_ext
        except:
            return "❓ Cannot process"

def extract_text_from_pptx(filepath):
    """Extract text content and keywords from PowerPoint document"""
    try:
        # First check if file exists and is readable
        if not os.path.exists(filepath):
            logger.warning(f"PowerPoint document not found: {filepath}")
            # Fallback to filename extraction
            filename = os.path.basename(filepath)
            name_without_ext = os.path.splitext(filename)[0]
            result = extract_keywords(name_without_ext)
            return result if result else "📄 " + os.path.splitext(filename)[0]
            
        try:
            prs = Presentation(filepath)

            # Extract slide titles and content
            slide_texts = []

            for slide in prs.slides:
                slide_text = ""

                # Extract title (usually the first shape)
                title = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip() and not title:
                        title = shape.text.strip()
                        break

                # Extract all text
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())

                # Combine text
                if title:
                    slide_text = f"{title}: {' '.join(texts)}"
                else:
                    slide_text = " ".join(texts)

                if slide_text:
                    slide_texts.append(slide_text)

                # Only process first 3 slides
                if len(slide_texts) >= 3:
                    break

            # Combine all slide text
            full_text = " ".join(slide_texts)

            # Extract keywords
            result = extract_keywords(full_text)
            
            # If no content extracted, fallback to filename
            if not result or result == "No keywords":
                logger.warning(f"No content extracted from PowerPoint, using filename: {filepath}")
                filename = os.path.basename(filepath)
                name_without_ext = os.path.splitext(filename)[0]
                result = extract_keywords(name_without_ext)
                
            return result if result else "📄 " + os.path.splitext(os.path.basename(filepath))[0]
            
        except Exception as ppt_error:
            # Log specific PowerPoint parsing error
            logger.warning(f"Failed to parse PowerPoint content: {filepath} - {str(ppt_error)}")
            # Fallback to filename extraction
            filename = os.path.basename(filepath)
            name_without_ext = os.path.splitext(filename)[0]
            result = extract_keywords(name_without_ext)
            logger.info(f"Using filename keywords for {filepath}: {result}")
            return result if result else "📄 " + name_without_ext
            
    except Exception as e:
        # Final fallback for any other errors
        logger.warning(f"Error processing PowerPoint {filepath}, using filename: {str(e)}")
        try:
            filename = os.path.basename(filepath)
            name_without_ext = os.path.splitext(filename)[0]
            result = extract_keywords(name_without_ext)
            return result if result else "📄 " + name_without_ext
        except:
            return "❓ Cannot process"

def extract_text_from_xlsx(filepath):
    """Extract keywords from Excel file by using only the filename (no extraction or reading of file content)"""
    try:
        filename = os.path.basename(filepath)
        name_without_ext = os.path.splitext(filename)[0]
        result = extract_keywords(name_without_ext)
        return result if result else "❓ No content"
    except Exception as e:
        logger.error(f"Error classifying Excel file by filename: {str(e)}")
        return "❓ Cannot read"

def extract_text_from_dxf(filepath):
    """Extract keywords from DXF file by using only the filename (no extraction or reading of file content)"""
    try:
        filename = os.path.basename(filepath)
        name_without_ext = os.path.splitext(filename)[0]
        result = extract_keywords(name_without_ext)
        return result if result else "❓ No content"
    except Exception as e:
        logger.error(f"Error classifying DXF file by filename: {str(e)}")
        return "❓ Cannot read"

def extract_keywords(text, max_keywords=5):
    """Extract keywords from text"""
    if not text or not isinstance(text, str):
        return ""

    # Simple stop words list (can be expanded as needed)
    stop_words = set(['a', 'an', 'the', 'and', 'or', 'but', 'if', 'then', 'else', 'when',
                      'at', 'from', 'by', 'for', 'with', 'about', 'against', 'between',
                      'into', 'through', 'during', 'before', 'after', 'above', 'below',
                      'to', 'of', 'in', 'on', 'is', 'are', 'was', 'were', 'be', 'been',
                      'being', 'have', 'has', 'had', 'having', 'do', 'does', 'did', 'doing',
                      'the', 'and', 'or', 'is', 'in', 'on', 'has', 'has', 'had', 'having',
                      'do', 'does', 'did', 'doing'])

    # Split text into words
    words = text.lower().split()

    # Filter stop words and short words
    filtered_words = [word for word in words if word not in stop_words and len(word) > 1]

    # Calculate word frequency
    word_freq = {}
    for word in filtered_words:
        word_freq[word] = word_freq.get(word, 0) + 1

    # Sort by frequency and take top N
    sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
    top_keywords = [word for word, _ in sorted_words[:max_keywords]]

    return " ".join(top_keywords) if top_keywords else "No keywords"

def extract_text_from_zip(filepath):
    """Extract keywords from zip file by using only the zip filename (no extraction or reading of file content)"""
    try:
        filename = os.path.basename(filepath)
        name_without_ext = os.path.splitext(filename)[0]
        result = extract_keywords(name_without_ext)
        return result if result else "❓ No content"
    except Exception as e:
        logger.error(f"Error classifying ZIP file by filename: {str(e)}")
        return "❓ Cannot read"

########################Resume Upload Module#############################################################